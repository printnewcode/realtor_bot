from pptx import Presentation
from pptx.util import Inches
import os

def create_presentation():
    # Создаем объект презентации
    prs = Presentation()
    file_0=os.path.join("files", "file_0.png")
    file_1=os.path.join("files", "file_1.jpg")
    # Получаем количество слайдов от пользователя
    slide_count = int(input("Введите количество слайдов для презентации: "))
    
    for i in range(slide_count):
        # Создаем новый слайд
        slide_layout = prs.slide_layouts[5]  # пустой слайд
        slide = prs.slides.add_slide(slide_layout)

        # Ввод текста от пользователя
        text = input(f"Введите текст для слайда {i+1}: ")
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = text

        # Ввод изображения от пользователя
        img_path = file_0 if i == 0 else file_1
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1), Inches(2.5), width=Inches(8))
        else:
            print(f"Файл {img_path} не найден. Изображение не будет добавлено.")

    # Сохранение презентации
    prs_name = input("Введите имя файла для сохранения презентации (например, presentation.pptx): ")
    prs.save(prs_name)
    print(f"Презентация {prs_name} успешно создана!")
create_presentation()